import io
import logging
from fastapi import APIRouter, UploadFile, File, HTTPException

router = APIRouter()
logger = logging.getLogger(__name__)

TIPE_FILE_DIIZINKAN = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}

MAKS_UKURAN = 20 * 1024 * 1024


@router.post("/parse-file")
async def parse_file(file: UploadFile = File(...)):
    if file.content_type not in TIPE_FILE_DIIZINKAN:
        raise HTTPException(
            status_code=400,
            detail=f"Tipe file tidak didukung: {file.content_type}. Gunakan PDF, DOCX, atau PPTX.",
        )

    konten = await file.read()

    if len(konten) > MAKS_UKURAN:
        raise HTTPException(status_code=400, detail="Ukuran file melebihi batas 20MB")

    tipe = TIPE_FILE_DIIZINKAN[file.content_type]

    try:
        if tipe == "pdf":
            teks, jumlah_halaman = ekstrak_pdf(konten)
        elif tipe == "docx":
            teks, jumlah_halaman = ekstrak_docx(konten)
        elif tipe == "pptx":
            teks, jumlah_halaman = ekstrak_pptx(konten)
        else:
            raise HTTPException(status_code=400, detail="Tipe tidak dikenali")

        teks_bersih = bersihkan_teks(teks)

        return {
            "teks_hasil": teks_bersih,
            "jumlah_halaman": jumlah_halaman,
            "tipe_file": tipe,
            "ukuran_kb": len(konten) // 1024,
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Gagal parse file: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Gagal memproses file: {str(e)}")


def ekstrak_pdf(konten: bytes) -> tuple[str, int]:
    import fitz

    doc = fitz.open(stream=konten, filetype="pdf")
    teks_semua = []
    for halaman in doc:
        teks_semua.append(halaman.get_text())
    jumlah = len(doc)
    doc.close()
    return "\n\n".join(teks_semua), jumlah


def ekstrak_docx(konten: bytes) -> tuple[str, int]:
    from docx import Document

    doc = Document(io.BytesIO(konten))
    paragraf = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraf), len(paragraf) // 25 + 1


def ekstrak_pptx(konten: bytes) -> tuple[str, int]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(konten))
    teks_semua = []
    for slide in prs.slides:
        teks_slide = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                teks_slide.append(shape.text)
        if teks_slide:
            teks_semua.append("\n".join(teks_slide))
    return "\n\n".join(teks_semua), len(prs.slides)


def bersihkan_teks(teks: str) -> str:
    import re

    teks = re.sub(r"\n{3,}", "\n\n", teks)
    teks = re.sub(r"[ \t]+", " ", teks)
    teks = "\n".join(line.strip() for line in teks.split("\n"))
    return teks.strip()
